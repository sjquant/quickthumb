"""PPTX (PowerPoint) exporter built on python-pptx.

The canvas becomes a single slide sized to the canvas pixel dimensions.
Backgrounds, outlines, shapes, and text map to native, editable PowerPoint
shapes and text boxes; raster images, SVG overlays, blend modes, and custom
layers are embedded as pixel-exact PNG pictures rendered by the PIL pipeline.

PowerPoint lays text out with its own font metrics, so text placement is a
close approximation rather than pixel-identical to the PNG render.
"""

from __future__ import annotations

import contextlib
import math
from io import BytesIO
from typing import TYPE_CHECKING, cast

from PIL import ImageFont

from quickthumb._base import (
    apply_alignment,
    expanded_rotation_size,
    parse_coordinate,
    parse_padding,
)
from quickthumb._export_base import (
    Box,
    RasterFragment,
    TextBlockLayout,
    TextRunLayout,
    color_to_rgba,
    compute_text_layout,
    flatten_layers,
    rasterize_layers,
    split_backdrop_prefix,
    uses_image_fill,
    uses_radial_fill,
)
from quickthumb.errors import RenderingError
from quickthumb.models import (
    Animation,
    BackgroundLayer,
    Blinds,
    Checkerboard,
    Glow,
    LinearGradient,
    OutlineLayer,
    RadialGradient,
    Shadow,
    ShapeLayer,
    Stroke,
    TextLayer,
    Transition,
    Wheel,
    Wipe,
)
from quickthumb.models import Box as BoxAnimation  # avoid clash with _export_base.Box

if TYPE_CHECKING:
    from quickthumb.canvas import Canvas, RenderableLayer

EMU_PER_PX = 9525  # 96 dpi
PT_PER_PX = 0.75
PPTX_MIN_SLIDE_EMU = 914400  # 1 inch
PPTX_MAX_SLIDE_EMU = 51206400  # 56 inches

_STAR_SHAPES = {4, 5, 6, 7, 8, 10, 12, 16, 24, 32}

_ALIGN_MAP = {"left": "LEFT", "center": "CENTER", "right": "RIGHT"}

# 2010 transition extension namespace, used for a precise transition duration in
# milliseconds (the base spd attribute only offers slow/med/fast).
_P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"

# presetID values PowerPoint uses to label entrance/exit effects in its UI. The
# actual motion comes from the behaviour XML, so an approximate id is harmless.
_ANIM_PRESET_IDS = {
    "appear": 1,
    "fade": 10,
    "wipe": 22,
    "blinds": 3,
    "checkerboard": 5,
    "circle": 6,
    "diamond": 7,
    "dissolve": 9,
    "wheel": 21,
    "box": 8,
}


def _emu(px: float) -> int:
    return round(px * EMU_PER_PX)


def _require_pptx():
    try:
        import pptx  # noqa: F401
    except ImportError:
        raise RenderingError(
            "python-pptx is required for PPTX export. "
            "Install it with: pip install 'quickthumb[pptx]'"
        ) from None


class PptxExporter:
    def __init__(self, canvas: Canvas | None = None):
        _require_pptx()
        # canvas is the single-slide convenience target; the multi-slide paths
        # take their canvases explicitly and leave this None.
        self._canvas = canvas

    def export_bytes(self) -> bytes:
        buffer = BytesIO()
        self.save(buffer)
        return buffer.getvalue()

    def save(self, path_or_stream) -> None:
        if self._canvas is None:
            raise RenderingError("PptxExporter() needs a canvas for single-slide export.")
        presentation = self._build([self._canvas])
        presentation.save(path_or_stream)

    def save_canvases(
        self,
        canvases: list[Canvas],
        path_or_stream,
        transitions: list[Transition | None] | None = None,
    ) -> None:
        """Write one or more canvases to a multi-slide presentation (one slide per canvas)."""
        presentation = self._build(canvases, transitions=transitions)
        presentation.save(path_or_stream)

    def export_bytes_canvases(
        self, canvases: list[Canvas], transitions: list[Transition | None] | None = None
    ) -> bytes:
        buffer = BytesIO()
        self.save_canvases(canvases, buffer, transitions=transitions)
        return buffer.getvalue()

    def _build(
        self, canvases: list[Canvas], transitions: list[Transition | None] | None = None
    ):
        from pptx import Presentation
        from pptx.util import Emu

        if not canvases:
            raise RenderingError("Cannot export a presentation with no slides.")

        presentation = Presentation()
        # The presentation has a single page size; the first slide defines it.
        # Validate before sizing so out-of-range canvases raise our error, not
        # python-pptx's.
        first = canvases[0]
        self._canvas = first
        self._validate_slide_dimensions()
        presentation.slide_width = Emu(_emu(first.width))
        presentation.slide_height = Emu(_emu(first.height))

        for index, canvas in enumerate(canvases):
            canvas._validate_image_paths()
            canvas._ctx.begin_render_pass()
            self._canvas = canvas
            if canvas is not first:  # the first slide was already validated above
                self._validate_slide_dimensions()
            self._slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            # (animation-or-list, [shape_id, ...]) records collected per layer.
            self._anim_records: list[tuple[object, list[int]]] = []

            prefix, rest = split_backdrop_prefix(flatten_layers(canvas))
            if prefix:
                fragment = rasterize_layers(canvas, prefix)
                if fragment:
                    self._add_fragment(fragment)
            for layer in rest:
                self._emit_tracked_layer(layer)

            transition = transitions[index] if transitions else None
            if transition is not None:
                self._apply_transition(transition)
            entries = self._resolve_animation_entries()
            if entries:
                self._apply_timing(entries)

        return presentation

    def _validate_slide_dimensions(self) -> None:
        canvas = self._canvas
        width_emu = _emu(canvas.width)
        height_emu = _emu(canvas.height)
        if (
            PPTX_MIN_SLIDE_EMU <= width_emu <= PPTX_MAX_SLIDE_EMU
            and PPTX_MIN_SLIDE_EMU <= height_emu <= PPTX_MAX_SLIDE_EMU
        ):
            return

        min_px = math.ceil(PPTX_MIN_SLIDE_EMU / EMU_PER_PX)
        max_px = math.floor(PPTX_MAX_SLIDE_EMU / EMU_PER_PX)
        raise RenderingError(
            "PPTX export supports canvas dimensions between "
            f"{min_px}px and {max_px}px in each direction."
        )

    # ------------------------------------------------------------------ layers

    def _emit_tracked_layer(self, layer: RenderableLayer):
        """Emit a layer, recording the shape ids it produces if it is animated.

        Animations target shapes by id, so we diff the slide's shapes before and
        after emitting; whatever is new belongs to this layer. A single layer can
        produce several shapes (e.g. a text box plus its background fills), and
        the animation is applied to all of them together.
        """
        animation = getattr(layer, "animation", None)
        if animation is None:
            self._emit_layer(layer)
            return

        before = {shape.shape_id for shape in self._slide.shapes}
        self._emit_layer(layer)
        new_ids = [
            shape.shape_id for shape in self._slide.shapes if shape.shape_id not in before
        ]
        if new_ids:
            self._anim_records.append((animation, new_ids))

    def _resolve_animation_entries(self) -> list[tuple[Animation, list[int]]]:
        """Turn the per-layer records into (animation, shape_ids) effect entries.

        Consecutive records that share the *same* animation object are the
        flattened children of one animated group; they merge into a single target
        set so the group animates as one effect. Each record's animation (a single
        effect or a list) then expands into one entry per effect.
        """
        merged: list[tuple[object, list[int]]] = []
        for animation, ids in self._anim_records:
            if merged and merged[-1][0] is animation:
                merged[-1][1].extend(ids)
            else:
                merged.append((animation, list(ids)))

        entries: list[tuple[Animation, list[int]]] = []
        for animation, ids in merged:
            anims = animation if isinstance(animation, list) else [animation]
            for anim in anims:
                entries.append((anim, ids))
        return entries

    def _emit_layer(self, layer: RenderableLayer):
        if isinstance(layer, BackgroundLayer):
            self._emit_background(layer)
        elif isinstance(layer, OutlineLayer):
            self._emit_outline(layer)
        elif isinstance(layer, ShapeLayer):
            self._emit_shape(layer)
        elif isinstance(layer, TextLayer):
            self._emit_text(layer)
        else:
            self._emit_raster_fallback(layer)

    def _emit_raster_fallback(self, layer: RenderableLayer):
        fragment = rasterize_layers(self._canvas, [layer])
        if fragment:
            self._add_fragment(fragment)

    def _add_fragment(self, fragment: RasterFragment):
        from pptx.util import Emu

        self._slide.shapes.add_picture(
            BytesIO(fragment.png_bytes),
            Emu(_emu(fragment.x)),
            Emu(_emu(fragment.y)),
            Emu(_emu(fragment.width)),
            Emu(_emu(fragment.height)),
        )

    # -------------------------------------------------------------- background

    def _emit_background(self, layer: BackgroundLayer):
        if layer.effects or layer.image or isinstance(layer.gradient, RadialGradient):
            self._emit_raster_fallback(layer)
            return

        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        canvas = self._canvas
        shape = self._slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(_emu(canvas.width)), Emu(_emu(canvas.height))
        )
        shape.line.fill.background()
        shape.shadow.inherit = False

        if layer.color is not None:
            rgba = color_to_rgba(canvas, layer.color, layer.opacity)
            self._set_solid_fill(shape, rgba)
        elif isinstance(layer.gradient, LinearGradient):
            stops = self._remap_linear_stops(
                layer.gradient.stops, layer.gradient.angle, canvas.width, canvas.height
            )
            self._set_gradient_fill_xml(
                shape._element.spPr, stops, layer.gradient.angle, layer.opacity
            )

    def _remap_linear_stops(
        self, stops: list[tuple[str, float]], angle: float, width: float, height: float
    ) -> list[tuple[tuple, float]]:
        """Remap stop positions from PIL's diagonal-spanning ramp to a box-spanning ramp.

        PowerPoint stretches a linear gradient across the shape's projection onto
        the gradient axis; PIL stretches it across the full diagonal, centered.
        """
        parsed = sorted(
            ((color_to_rgba(self._canvas, color), position) for color, position in stops),
            key=lambda stop: stop[1],
        )
        diagonal = math.ceil(math.hypot(width, height))
        theta = math.radians(angle)
        extent = abs(math.cos(theta)) * width + abs(math.sin(theta)) * height
        if extent <= 0:
            return parsed

        def to_box(position: float) -> float:
            return (position * diagonal - diagonal / 2 + extent / 2) / extent

        def color_at(position: float) -> tuple:
            if position <= parsed[0][1]:
                return parsed[0][0]
            if position >= parsed[-1][1]:
                return parsed[-1][0]
            for (c1, p1), (c2, p2) in zip(parsed, parsed[1:], strict=False):
                if p1 <= position <= p2:
                    ratio = (position - p1) / (p2 - p1) if p2 != p1 else 0
                    return tuple(int(a + (b - a) * ratio) for a, b in zip(c1, c2, strict=True))
            return parsed[-1][0]

        boundary_lo = (diagonal - extent) / (2 * diagonal)
        boundary_hi = (diagonal + extent) / (2 * diagonal)
        remapped: list[tuple[tuple, float]] = [(color_at(boundary_lo), 0.0)]
        for color, position in parsed:
            q = to_box(position)
            if 0.0 < q < 1.0:
                remapped.append((color, q))
        remapped.append((color_at(boundary_hi), 1.0))
        return remapped

    # ----------------------------------------------------------------- outline

    def _emit_outline(self, layer: OutlineLayer):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        canvas = self._canvas
        inset = layer.offset + layer.width / 2
        shape = self._slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(_emu(inset)),
            Emu(_emu(inset)),
            Emu(_emu(canvas.width - 2 * layer.offset - layer.width)),
            Emu(_emu(canvas.height - 2 * layer.offset - layer.width)),
        )
        shape.fill.background()
        shape.shadow.inherit = False
        rgba = color_to_rgba(canvas, layer.color, layer.opacity)
        self._set_line(shape, rgba, layer.width)

    # ------------------------------------------------------------------ shapes

    def _emit_shape(self, layer: ShapeLayer):
        canvas = self._canvas
        x = parse_coordinate(layer.position[0], canvas.width)
        y = parse_coordinate(layer.position[1], canvas.height)
        width, height = layer.width, layer.height

        expanded = expanded_rotation_size((width, height), layer.rotation)
        paste_x, paste_y = (x, y)
        if layer.align:
            paste_x, paste_y = apply_alignment(x, y, expanded, layer.align)
        left = paste_x + (expanded[0] - width) / 2
        top = paste_y + (expanded[1] - height) / 2

        shape = self._add_autoshape(layer, left, top)
        if shape is None:
            self._emit_raster_fallback(layer)
            return

        shape.shadow.inherit = False
        if layer.rotation != 0:
            shape.rotation = layer.rotation

        rgba = color_to_rgba(canvas, layer.color, layer.opacity)
        self._set_solid_fill(shape, rgba)

        strokes = [e for e in layer.effects if isinstance(e, Stroke)]
        if strokes:
            stroke = strokes[0]
            self._set_line(shape, color_to_rgba(canvas, stroke.color, layer.opacity), stroke.width)
        else:
            from pptx.shapes.autoshape import Shape

            cast(Shape, shape).line.fill.background()

        glows = [e for e in layer.effects if isinstance(e, Glow)]
        shadows = [e for e in layer.effects if isinstance(e, Shadow)]
        if glows or shadows:
            self._append_effects(shape._element.spPr, glows, shadows, opacity=layer.opacity)

    def _add_autoshape(self, layer: ShapeLayer, left: float, top: float):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        width, height = layer.width, layer.height
        args = (Emu(_emu(left)), Emu(_emu(top)), Emu(_emu(width)), Emu(_emu(height)))
        shapes = self._slide.shapes

        if layer.shape == "rectangle":
            if layer.border_radius > 0:
                shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, *args)
                ratio = min(0.5, layer.border_radius / min(width, height))
                self._set_adjustment(shape, ratio)
            else:
                shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, *args)
            return shape
        if layer.shape == "ellipse":
            return shapes.add_shape(MSO_SHAPE.OVAL, *args)
        if layer.shape == "pill":
            shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, *args)
            self._set_adjustment(shape, 0.5)
            return shape
        if layer.shape == "triangle":
            return shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, *args)
        if layer.shape == "star" and layer.star_points in _STAR_SHAPES:
            mso_star = getattr(MSO_SHAPE, f"STAR_{layer.star_points}_POINT")
            shape = shapes.add_shape(mso_star, *args)
            # DrawingML star adjustment is inner-radius/diameter (max 0.5).
            self._set_adjustment(shape, layer.inner_radius * 0.5)
            return shape
        return self._add_freeform(layer, left, top)

    def _add_freeform(self, layer: ShapeLayer, left: float, top: float):
        from pptx.util import Emu

        normalized = self._canvas._shapes.normalized_shape_points(layer)
        if not normalized:
            return None
        points = [(px * layer.width, py * layer.height) for px, py in normalized]
        builder = self._slide.shapes.build_freeform(
            start_x=points[0][0], start_y=points[0][1], scale=EMU_PER_PX
        )
        builder.add_line_segments(points[1:], close=True)
        return builder.convert_to_shape(origin_x=Emu(_emu(left)), origin_y=Emu(_emu(top)))

    @staticmethod
    def _set_adjustment(shape, value: float):
        # Shapes without adjustment handles keep their default geometry.
        with contextlib.suppress(IndexError, ValueError):
            shape.adjustments[0] = value

    # -------------------------------------------------------------------- text

    def _emit_text(self, layer: TextLayer):
        if uses_image_fill(layer) or uses_radial_fill(layer):
            self._emit_raster_fallback(layer)
            return

        layout = compute_text_layout(self._canvas, layer)
        if not any(layout.lines) or layout.block_box is None:
            return

        if layout.block_backgrounds and layout.block_bg_box:
            for background in layout.block_backgrounds:
                self._add_text_background(background, layout.block_bg_box, layout)
        for line in layout.lines:
            for run in line:
                for background in run.backgrounds:
                    if run.bg_box:
                        self._add_text_background(background, run.bg_box, layout)

        self._add_textbox(layer, layout)

    def _add_text_background(self, background, content_box: Box, layout: TextBlockLayout):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        pad_top, pad_right, pad_bottom, pad_left = parse_padding(background.padding)
        box = (
            content_box[0] - pad_left,
            content_box[1] - pad_top,
            content_box[2] + pad_right,
            content_box[3] + pad_bottom,
        )
        left, top, width, height, rotation = self._place_box(layout, box)
        shape_type = (
            MSO_SHAPE.ROUNDED_RECTANGLE if background.border_radius > 0 else MSO_SHAPE.RECTANGLE
        )
        shape = self._slide.shapes.add_shape(
            shape_type, Emu(_emu(left)), Emu(_emu(top)), Emu(_emu(width)), Emu(_emu(height))
        )
        if background.border_radius > 0:
            self._set_adjustment(shape, min(0.5, background.border_radius / min(width, height)))
        shape.line.fill.background()
        shape.shadow.inherit = False
        if rotation:
            shape.rotation = rotation
        rgba = color_to_rgba(self._canvas, background.color, background.opacity * layout.opacity)
        self._set_solid_fill(shape, rgba)

    def _add_textbox(self, layer: TextLayer, layout: TextBlockLayout):
        from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Emu, Pt

        assert layout.block_box is not None  # guarded by _emit_text
        left, top, width, height, rotation = self._place_box(layout, layout.block_box)
        textbox = self._slide.shapes.add_textbox(
            Emu(_emu(left)), Emu(_emu(top)), Emu(_emu(width)), Emu(_emu(height))
        )
        if rotation:
            textbox.rotation = rotation

        frame = textbox.text_frame
        frame.word_wrap = False
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.vertical_anchor = MSO_ANCHOR.TOP
        frame.margin_left = frame.margin_right = Emu(0)
        frame.margin_top = frame.margin_bottom = Emu(0)

        horizontal = layer.align.horizontal if layer.align else "left"
        alignment = getattr(PP_ALIGN, _ALIGN_MAP.get(horizontal, "LEFT"))

        for index, line in enumerate(layout.lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = alignment
            if index < len(layout.line_heights):
                paragraph.line_spacing = Pt(layout.line_heights[index] * PT_PER_PX)
            for run_layout in line:
                self._add_run(paragraph, run_layout, layout.opacity)

    def _add_run(self, paragraph, run_layout: TextRunLayout, layer_opacity: float):
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        run = paragraph.add_run()
        run.text = run_layout.text
        font = run.font
        font.size = Pt(run_layout.size * PT_PER_PX)
        font.bold = run_layout.bold or self._weight_is_bold(run_layout.weight)
        font.italic = run_layout.italic
        if isinstance(run_layout.font, ImageFont.FreeTypeFont):
            font.name = run_layout.font.getname()[0]

        rpr = run._r.get_or_add_rPr()  # noqa: N802 - python-pptx API name
        if isinstance(run_layout.fill, (LinearGradient, RadialGradient)) and run_layout.fill_box:
            box = run_layout.fill_box
            if isinstance(run_layout.fill, LinearGradient):
                stops = self._remap_linear_stops(
                    run_layout.fill.stops,
                    run_layout.fill.angle,
                    box[2] - box[0],
                    box[3] - box[1],
                )
                self._set_gradient_fill_xml(rpr, stops, run_layout.fill.angle, layer_opacity)
            else:
                first_color = color_to_rgba(
                    self._canvas, sorted(run_layout.fill.stops, key=lambda s: s[1])[0][0]
                )
                font.color.rgb = RGBColor(*first_color[:3])
        else:
            rgba = run_layout.color
            font.color.rgb = RGBColor(*rgba[:3])
            alpha = (rgba[3] / 255) * layer_opacity
            if alpha < 1:
                self._append_alpha(rpr, alpha)

        if run_layout.strokes:
            stroke = run_layout.strokes[0]
            self._prepend_text_outline(
                rpr,
                color_to_rgba(self._canvas, stroke.color, layer_opacity),
                stroke.width,
            )
        if run_layout.glows or run_layout.shadows:
            self._append_effects(rpr, run_layout.glows, run_layout.shadows, opacity=layer_opacity)

    @staticmethod
    def _weight_is_bold(weight: int | str | None) -> bool:
        if isinstance(weight, int):
            return weight >= 600
        if isinstance(weight, str):
            return weight.lower() in ("bold", "semibold", "extrabold", "black") or (
                weight.isdigit() and int(weight) >= 600
            )
        return False

    def _place_box(
        self, layout: TextBlockLayout, box: Box
    ) -> tuple[float, float, float, float, float]:
        """Map a (possibly rotation-local) box to slide coordinates and a rotation."""
        width, height = box[2] - box[0], box[3] - box[1]
        if layout.rotation == 0 or layout.origin is None or layout.rot_center_local is None:
            return box[0], box[1], width, height, 0.0

        center_x = (box[0] + box[2]) / 2 - layout.rot_center_local[0]
        center_y = (box[1] + box[3]) / 2 - layout.rot_center_local[1]
        theta = math.radians(layout.rotation)
        cos_t, sin_t = math.cos(theta), math.sin(theta)
        rotated_x = center_x * cos_t - center_y * sin_t
        rotated_y = center_x * sin_t + center_y * cos_t
        global_cx = layout.origin[0] + layout.rot_center_local[0] + rotated_x
        global_cy = layout.origin[1] + layout.rot_center_local[1] + rotated_y
        return global_cx - width / 2, global_cy - height / 2, width, height, layout.rotation

    # --------------------------------------------------------------- XML hacks

    def _set_solid_fill(self, shape, rgba: tuple):
        from pptx.dml.color import RGBColor

        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgba[:3])
        if rgba[3] < 255:
            self._append_alpha(shape._element.spPr, rgba[3] / 255)

    def _set_line(self, shape, rgba: tuple, width_px: float):
        from pptx.dml.color import RGBColor
        from pptx.util import Emu

        shape.line.color.rgb = RGBColor(*rgba[:3])
        shape.line.width = Emu(_emu(width_px))
        if rgba[3] < 255:
            ln = shape.line._get_or_add_ln()
            self._append_alpha(ln, rgba[3] / 255)

    @staticmethod
    def _append_alpha(parent_element, alpha: float):
        """Append an <a:alpha> child to the first srgbClr under parent_element."""
        from pptx.oxml.ns import qn

        srgb = parent_element.find(f".//{qn('a:srgbClr')}")
        if srgb is None:
            return
        alpha_element = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
        srgb.append(alpha_element)

    @staticmethod
    def _srgb_xml(rgba: tuple) -> str:
        alpha = ""
        if rgba[3] < 255:
            alpha = f'<a:alpha val="{int(rgba[3] / 255 * 100000)}"/>'
        return '<a:srgbClr val="{:02X}{:02X}{:02X}">{}</a:srgbClr>'.format(*rgba[:3], alpha)

    def _set_gradient_fill_xml(
        self, properties_element, stops: list[tuple[tuple, float]], angle: float, opacity: float
    ):
        """Replace the fill of an spPr/rpr element with a multi-stop linear gradient."""
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsdecls, qn

        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
            for element in properties_element.findall(qn(tag)):
                properties_element.remove(element)

        gradient_stops = []
        for rgba, position in stops:
            if opacity < 1:
                rgba = (*rgba[:3], int(rgba[3] * opacity))
            gradient_stops.append(
                f'<a:gs pos="{int(round(max(0.0, min(1.0, position)) * 100000))}">'
                f"{self._srgb_xml(rgba)}</a:gs>"
            )
        ang = int(round((angle % 360) * 60000))
        xml = (
            f'<a:gradFill {nsdecls("a")} rotWithShape="1"><a:gsLst>'
            + "".join(gradient_stops)
            + f'</a:gsLst><a:lin ang="{ang}" scaled="0"/></a:gradFill>'
        )
        properties_element.insert_element_before(
            parse_xml(xml),
            "a:ln",
            "a:effectLst",
            "a:effectDag",
            "a:scene3d",
            "a:sp3d",
            "a:extLst",
            "a:highlight",
            "a:uLnTx",
            "a:uLn",
            "a:uFillTx",
            "a:uFill",
            "a:latin",
        )

    def _append_effects(
        self,
        properties_element,
        glows: list[Glow],
        shadows: list[Shadow],
        opacity: float = 1.0,
    ):
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsdecls, qn

        existing = properties_element.find(qn("a:effectLst"))
        if existing is not None:
            properties_element.remove(existing)

        parts = []
        for glow in glows:
            rgba = color_to_rgba(self._canvas, glow.color, glow.opacity * opacity)
            parts.append(f'<a:glow rad="{_emu(glow.radius * 3)}">{self._srgb_xml(rgba)}</a:glow>')
        for shadow in shadows:
            rgba = color_to_rgba(self._canvas, shadow.color, opacity)
            distance = math.hypot(shadow.offset_x, shadow.offset_y)
            direction = int(
                round(math.degrees(math.atan2(shadow.offset_y, shadow.offset_x)) % 360 * 60000)
            )
            parts.append(
                f'<a:outerShdw blurRad="{_emu(shadow.blur_radius * 2)}" '
                f'dist="{_emu(distance)}" dir="{direction}" rotWithShape="0">'
                f"{self._srgb_xml(rgba)}</a:outerShdw>"
            )
        xml = f"<a:effectLst {nsdecls('a')}>" + "".join(parts) + "</a:effectLst>"
        properties_element.insert_element_before(
            parse_xml(xml),
            "a:effectDag",
            "a:scene3d",
            "a:sp3d",
            "a:extLst",
            "a:highlight",
            "a:uLnTx",
            "a:uLn",
            "a:uFillTx",
            "a:uFill",
            "a:latin",
        )

    def _prepend_text_outline(self, rpr, rgba: tuple, width_px: float):
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsdecls, qn

        existing = rpr.find(qn("a:ln"))
        if existing is not None:
            rpr.remove(existing)
        xml = (
            f'<a:ln {nsdecls("a")} w="{_emu(width_px)}"><a:solidFill>'
            f"{self._srgb_xml(rgba)}</a:solidFill></a:ln>"
        )
        rpr.insert(0, parse_xml(xml))

    # ------------------------------------------------------ slide transitions

    def _apply_transition(self, transition: Transition):
        """Append a <p:transition> element describing this slide's transition."""
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsdecls

        child = self._transition_child(transition.effect, transition.direction)
        duration = transition.duration
        speed = "fast" if duration <= 0.5 else "med" if duration <= 1.0 else "slow"

        attrs = [f'spd="{speed}"', f'advClick="{1 if transition.advance_on_click else 0}"']
        if transition.advance_after is not None:
            attrs.append(f'advTm="{int(round(transition.advance_after * 1000))}"')

        # spd alone is coarse (slow/med/fast); p14:dur carries the exact duration
        # for PowerPoint 2010+ and is declared inline so the element stays valid.
        xml = (
            f'<p:transition {nsdecls("p")} xmlns:p14="{_P14_NS}" '
            f'{" ".join(attrs)} p14:dur="{int(round(duration * 1000))}">'
            f"{child}</p:transition>"
        )
        self._slide._element.append(parse_xml(xml))

    def _transition_child(self, effect: str, direction: str | None) -> str:
        if effect == "none":
            return ""
        simple = {
            "cut": "cut",
            "fade": "fade",
            "dissolve": "dissolve",
            "newsflash": "newsflash",
            "wedge": "wedge",
            "wheel": "wheel",
            "circle": "circle",
            "diamond": "diamond",
            "random": "random",
        }
        if effect in simple:
            return f"<p:{simple[effect]}/>"
        if effect == "push":
            return f'<p:push dir="{self._dir_4(direction, "l")}"/>'
        if effect == "wipe":
            return f'<p:wipe dir="{self._dir_4(direction, "u")}"/>'
        if effect == "cover":
            return f'<p:cover dir="{self._dir_4(direction, "d")}"/>'
        if effect == "uncover":
            return f'<p:pull dir="{self._dir_4(direction, "d")}"/>'
        if effect == "zoom":
            return f'<p:zoom dir="{"out" if direction == "out" else "in"}"/>'
        if effect == "split":
            orient = "vert" if direction in ("vertical", "up", "down") else "horz"
            return f'<p:split orient="{orient}" dir="{"in" if direction == "in" else "out"}"/>'
        if effect in ("blinds", "checker", "comb"):
            return f'<p:{effect} dir="{self._dir_orient(direction)}"/>'
        return ""

    @staticmethod
    def _dir_4(direction: str | None, default: str) -> str:
        return {"left": "l", "right": "r", "up": "u", "down": "d"}.get(direction or "", default)

    @staticmethod
    def _dir_orient(direction: str | None) -> str:
        return "vert" if direction in ("vertical", "up", "down") else "horz"

    # ------------------------------------------------------- element animations

    def _apply_timing(self, animations: list[tuple[Animation, list[int]]]):
        """Append a <p:timing> tree animating the collected shapes in order."""
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsdecls

        # cTn ids 1 (tmRoot) and 2 (mainSeq) are fixed; the rest count up from 3.
        self._tn_id = 2
        groups = self._group_animations(animations)
        group_xml = "".join(self._build_click_group(group) for group in groups)

        seq = (
            '<p:seq concurrent="1" nextAc="seek">'
            '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f"{group_xml}"
            "</p:childTnLst></p:cTn>"
            '<p:prevCondLst><p:cond evt="onPrev" delay="0">'
            "<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>"
            '<p:nextCondLst><p:cond evt="onNext" delay="0">'
            "<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>"
            "</p:seq>"
        )
        xml = (
            f"<p:timing {nsdecls('p')}><p:tnLst>"
            '<p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            f"<p:childTnLst>{seq}</p:childTnLst></p:cTn></p:par>"
            "</p:tnLst></p:timing>"
        )
        self._slide._element.append(parse_xml(xml))

    @staticmethod
    def _group_animations(
        animations: list[tuple[Animation, list[int]]],
    ) -> list[list[tuple[Animation, list[int]]]]:
        """Split animations into click groups; with_previous joins the open group."""
        groups: list[list[tuple[Animation, list[int]]]] = []
        for anim, ids in animations:
            if not groups or anim.trigger != "with_previous":
                groups.append([])
            groups[-1].append((anim, ids))
        return groups

    def _build_click_group(self, group: list[tuple[Animation, list[int]]]) -> str:
        # The first effect's trigger decides whether the group waits for a click
        # (indefinite) or starts automatically after the previous one (0).
        outer_delay = "0" if group[0][0].trigger == "after_previous" else "indefinite"
        outer_id = self._next_id()
        inner_id = self._next_id()
        effects = "".join(
            self._build_effect_par(anim, ids, index)
            for index, (anim, ids) in enumerate(group)
        )
        return (
            "<p:par>"
            f'<p:cTn id="{outer_id}" fill="hold">'
            f'<p:stCondLst><p:cond delay="{outer_delay}"/></p:stCondLst>'
            "<p:childTnLst><p:par>"
            f'<p:cTn id="{inner_id}" fill="hold">'
            '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f"<p:childTnLst>{effects}</p:childTnLst>"
            "</p:cTn></p:par></p:childTnLst>"
            "</p:cTn></p:par>"
        )

    def _build_effect_par(self, anim: Animation, ids: list[int], index: int) -> str:
        if index == 0:
            node_type = "afterEffect" if anim.trigger == "after_previous" else "clickEffect"
        else:
            node_type = "withEffect"
        effect_id = self._next_id()
        preset_id = _ANIM_PRESET_IDS.get(anim.effect, 10)
        preset_class = "exit" if anim.animate == "exit" else "entr"
        behaviors = "".join(self._build_behaviors(anim, sid) for sid in ids)
        return (
            "<p:par>"
            f'<p:cTn id="{effect_id}" presetID="{preset_id}" presetClass="{preset_class}" '
            f'presetSubtype="0" fill="hold" grpId="0" nodeType="{node_type}">'
            f'<p:stCondLst><p:cond delay="{int(round(anim.delay * 1000))}"/></p:stCondLst>'
            f"<p:childTnLst>{behaviors}</p:childTnLst>"
            "</p:cTn></p:par>"
        )

    def _build_behaviors(self, anim: Animation, sid: int) -> str:
        duration_ms = int(round(anim.duration * 1000))
        if anim.animate == "exit":
            parts = []
            if anim.effect != "appear":
                parts.append(
                    self._anim_effect_xml("out", self._anim_filter(anim), duration_ms, sid)
                )
            # Pin the shape hidden once it has animated out.
            parts.append(self._set_visibility_xml(sid, "hidden", delay=duration_ms))
            return "".join(parts)

        # Entrance: reveal the shape, then play the chosen filter.
        parts = [self._set_visibility_xml(sid, "visible", delay=0)]
        if anim.effect != "appear":
            parts.append(
                self._anim_effect_xml("in", self._anim_filter(anim), duration_ms, sid)
            )
        return "".join(parts)

    def _set_visibility_xml(self, sid: int, value: str, delay: int) -> str:
        cid = self._next_id()
        return (
            "<p:set><p:cBhvr>"
            f'<p:cTn id="{cid}" dur="1" fill="hold">'
            f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
            "<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>"
            f'</p:cBhvr><p:to><p:strVal val="{value}"/></p:to></p:set>'
        )

    def _anim_effect_xml(self, transition: str, filt: str, duration_ms: int, sid: int) -> str:
        cid = self._next_id()
        return (
            f'<p:animEffect transition="{transition}" filter="{filt}"><p:cBhvr>'
            f'<p:cTn id="{cid}" dur="{duration_ms}"/>'
            f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
            "</p:cBhvr></p:animEffect>"
        )

    @staticmethod
    def _anim_filter(anim: Animation) -> str:
        """Map an animation effect to its PowerPoint animEffect filter string."""
        if isinstance(anim, Wipe):
            return f"wipe({anim.direction})"
        if isinstance(anim, Blinds):
            return f"blinds({anim.orientation})"
        if isinstance(anim, Checkerboard):
            return f"checkerboard({anim.direction})"
        if isinstance(anim, Wheel):
            return f"wheel({anim.spokes})"
        if isinstance(anim, BoxAnimation):
            return f"box({anim.direction})"
        # fade, circle, diamond, dissolve map straight through.
        return anim.effect

    def _next_id(self) -> int:
        self._tn_id += 1
        return self._tn_id
