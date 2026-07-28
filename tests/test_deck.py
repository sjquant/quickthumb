"""Tests for Deck: multi-slide / multi-image collections of Canvas objects."""

import json
from io import BytesIO
from pathlib import Path
from typing import cast

import pytest
from PIL import Image
from quickthumb import Canvas, Deck
from quickthumb.errors import RenderingError, ValidationError
from quickthumb.models import BackgroundLayer
from quickthumb.transitions import Dissolve, Fade, Split, Wheel

from tests._optional import require_pypdfium2

pptx = pytest.importorskip("pptx", reason="python-pptx is required for some deck tests")
from pptx import Presentation  # noqa: E402


def make_slide(text: str, width: int = 1280, height: int = 720) -> Canvas:
    """Build a small, self-contained slide with a solid background and label."""
    return (
        Canvas(width, height)
        .background(color="#101820")
        .text(content=text, size=96, color="#FFFFFF", position=("50%", "50%"), align="center")
    )


class TestDeckComposition:
    """Building a deck and inspecting its slides."""

    def test_should_collect_slides_via_constructor_and_chaining(self):
        """The constructor seeds slides and slide() appends more, staying chainable."""
        # given
        first = make_slide("1")
        second = make_slide("2")
        third = make_slide("3")

        # when
        deck = Deck(slides=[first]).slide(second).slide(third)

        # then
        assert list(deck) == [first, second, third]
        assert len(deck) == 3
        assert deck[1] is second

    def test_should_reject_non_canvas_slides(self):
        """Adding something that is not a Canvas raises a validation error."""
        # given
        deck = Deck()

        # when / then
        with pytest.raises(ValidationError, match="must be Canvas"):
            deck.slide("not a canvas")  # type: ignore[arg-type]

    def test_should_not_mutate_deck_or_canvas_when_audio_is_invalid(self):
        """Rejecting slide audio leaves the Deck and unsized Canvas unchanged."""
        # given
        deck = Deck(1280, 720)
        slide = Canvas().background(color="#101820")

        # when
        with pytest.raises(ValidationError, match="greater than or equal to 0"):
            deck.slide(slide, audio={"path": "voice.wav", "volume": -1})

        # then
        assert len(deck) == 0
        assert not slide.has_size

    def test_should_reject_non_string_speaker_notes_before_adding_slide(self):
        """Invalid speaker notes leave the deck unchanged."""
        # given: an empty deck and a valid canvas
        deck = Deck()
        slide = make_slide("1")

        # when / then: structured notes are rejected without appending the slide
        with pytest.raises(ValidationError, match="notes"):
            deck.slide(slide, notes=["not", "text"])  # type: ignore[arg-type]
        assert len(deck) == 0

    def test_should_reject_rendering_an_empty_deck(self, tmp_path: Path):
        """An empty deck cannot be rendered to any format."""
        # given
        deck = Deck()

        # when / then
        with pytest.raises(RenderingError, match="no slides"):
            deck.render(str(tmp_path / "out.pdf"))


class TestDefaultSize:
    """A deck size lets slides be written as bare, unsized canvases."""

    def test_should_inject_deck_size_into_unsized_slides(self):
        """An unsized Canvas added to a sized deck inherits the deck dimensions."""
        # given a deck with a default size and bare canvases
        deck = (
            Deck(1280, 720)
            .slide(Canvas().background(color="#101820"))
            .slide(Canvas().background(color="#1A1A2E"))
        )

        # then both slides take the deck size
        assert [(c.width, c.height) for c in deck] == [(1280, 720), (1280, 720)]

    def test_should_keep_explicit_slide_size_over_deck_default(self):
        """An explicitly sized canvas keeps its own dimensions inside a sized deck."""
        # given
        deck = Deck(1280, 720).slide(Canvas(1080, 1080).background(color="#101820"))

        # then the explicit size wins
        assert (deck[0].width, deck[0].height) == (1080, 1080)

    def test_should_derive_default_size_from_aspect_ratio(self):
        """from_aspect_ratio sets the deck default that unsized slides inherit."""
        # given
        deck = Deck.from_aspect_ratio("16:9", 1280).slide(Canvas().background(color="#101820"))

        # then
        assert (deck[0].width, deck[0].height) == (1280, 720)

    def test_should_reject_a_malformed_aspect_ratio(self):
        """A malformed ratio raises the library's ValidationError, not a raw error."""
        # when / then
        with pytest.raises(ValidationError, match="Aspect ratio"):
            Deck.from_aspect_ratio("16-9", 1280)

    def test_should_reject_unsized_slide_without_a_deck_size(self):
        """An unsized slide needs either a deck size or its own size."""
        # given a deck with no default size
        deck = Deck()

        # when / then
        with pytest.raises(ValidationError, match="no default size"):
            deck.slide(Canvas().background(color="#101820"))

    def test_should_reject_rendering_an_unsized_canvas_directly(self, tmp_path: Path):
        """A bare Canvas cannot be rendered until it is given a size."""
        # given
        canvas = Canvas().background(color="#101820")

        # when / then
        with pytest.raises(ValidationError, match="no size yet"):
            canvas.render(str(tmp_path / "out.png"))


class TestRenderDispatch:
    """render() dispatches on the output extension and rejects bad combinations."""

    def test_should_reject_single_svg_output(self, tmp_path: Path):
        """A deck has no single-file SVG form and says so explicitly."""
        # given
        deck = Deck().slide(make_slide("1"))

        # when / then
        with pytest.raises(RenderingError, match="single .svg"):
            deck.render(str(tmp_path / "deck.svg"))

    def test_should_reject_unknown_extension(self, tmp_path: Path):
        """An unrecognized extension is rejected rather than guessed."""
        # given
        deck = Deck().slide(make_slide("1"))

        # when / then
        with pytest.raises(RenderingError, match="Unsupported deck output format"):
            deck.render(str(tmp_path / "deck.docx"))

    def test_should_reject_format_override_for_document_output(self, tmp_path: Path):
        """A raster format override is meaningless for document output."""
        # given
        deck = Deck().slide(make_slide("1"))

        # when / then
        with pytest.raises(RenderingError, match="format override"):
            deck.render(str(tmp_path / "deck.pdf"), format="PNG")


class TestRasterSequence:
    """Rendering a deck to per-slide raster images."""

    def test_should_write_zero_padded_numbered_sequence(self, tmp_path: Path):
        """Raster output writes one file per slide with a zero-padded index."""
        # given
        deck = Deck(slides=[make_slide("a"), make_slide("b"), make_slide("c")])
        output = tmp_path / "slides.png"

        # when
        written = deck.render(str(output))

        # then
        assert written == [
            str(tmp_path / "slides_01.png"),
            str(tmp_path / "slides_02.png"),
            str(tmp_path / "slides_03.png"),
        ]
        for path in written:
            assert Path(path).exists()
            assert Image.open(path).size == (1280, 720)

    def test_should_pass_quality_through_for_jpeg_sequence(self, tmp_path: Path):
        """Raster sequences honor the quality argument for JPEG output."""
        # given
        deck = Deck(slides=[make_slide("a"), make_slide("b")])
        output = tmp_path / "slides.jpg"

        # when
        written = deck.render(str(output), quality=50)

        # then
        assert len(written) == 2
        assert all(Path(path).exists() for path in written)

    def test_should_not_write_partial_sequence_when_a_slide_fails(self, tmp_path: Path):
        """A missing asset on a later slide fails before any file is written."""
        # given a deck whose second slide references a missing image
        broken = make_slide("b").image(str(tmp_path / "missing.png"), position=(0, 0))
        deck = Deck(slides=[make_slide("a"), broken, make_slide("c")])
        output = tmp_path / "slides.png"

        # when / then: rendering raises and leaves no partial output behind
        with pytest.raises(FileNotFoundError):
            deck.render(str(output))
        assert list(tmp_path.glob("slides_*.png")) == []


class TestDocumentExport:
    """Rendering a deck to multi-page PDF and multi-slide PPTX documents."""

    def test_should_render_one_pdf_page_per_slide(self, tmp_path: Path):
        """A PDF deck produces a document with a page for every slide."""
        # given
        pdfium = require_pypdfium2()
        deck = Deck(slides=[make_slide("1"), make_slide("2"), make_slide("3")])
        output = tmp_path / "deck.pdf"

        # when
        deck.render(str(output))

        # then
        document = pdfium.PdfDocument(str(output))
        assert len(document) == 3

    def test_should_render_one_pptx_slide_per_slide(self, tmp_path: Path):
        """A PPTX deck produces a presentation with a slide for every slide."""
        # given
        deck = Deck(slides=[make_slide("1"), make_slide("2")])
        output = tmp_path / "deck.pptx"

        # when
        deck.render(str(output))

        # then
        presentation = Presentation(str(output))
        assert len(presentation.slides) == 2

    def test_should_expose_pdf_bytes_with_a_page_per_slide(self):
        """to_pdf() returns multi-page PDF bytes."""
        # given
        pdfium = require_pypdfium2()
        deck = Deck(slides=[make_slide("1"), make_slide("2")])

        # when
        data = deck.to_pdf()

        # then
        assert len(pdfium.PdfDocument(BytesIO(data))) == 2

    def test_should_expose_pptx_bytes_with_a_slide_per_slide(self):
        """to_pptx() returns multi-slide presentation bytes."""
        # given
        deck = Deck(slides=[make_slide("1"), make_slide("2"), make_slide("3")])

        # when
        data = deck.to_pptx()

        # then
        assert len(Presentation(BytesIO(data)).slides) == 3

    def test_should_reject_quality_for_document_output(self, tmp_path: Path):
        """Document formats do not accept the raster-only quality argument."""
        # given
        deck = Deck().slide(make_slide("1"))

        # when / then
        with pytest.raises(RenderingError, match="Quality parameter"):
            deck.render(str(tmp_path / "deck.pdf"), quality=80)

    def test_should_keep_first_slide_size_for_mixed_pptx(self, tmp_path: Path):
        """With mixed sizes the PPTX page size comes from the first slide."""
        # given
        deck = Deck(slides=[make_slide("wide", 1280, 720), make_slide("square", 800, 800)])
        output = tmp_path / "deck.pptx"

        # when
        deck.render(str(output))

        # then
        presentation = Presentation(str(output))
        assert presentation.slide_width == 1280 * 9525
        assert presentation.slide_height == 720 * 9525

    def test_should_size_each_pdf_page_to_its_own_slide(self, tmp_path: Path):
        """Unlike PPTX, a mixed-size PDF sizes every page to its slide."""
        # given a deck with two differently shaped slides
        pdfium = require_pypdfium2()
        deck = Deck(slides=[make_slide("wide", 1280, 720), make_slide("square", 800, 800)])
        output = tmp_path / "deck.pdf"

        # when
        deck.render(str(output))

        # then each page keeps its slide's pixel dimensions (1pt per pixel)
        document = pdfium.PdfDocument(str(output))
        assert [tuple(round(v) for v in document[i].get_size()) for i in range(len(document))] == [
            (1280, 720),
            (800, 800),
        ]


class TestDiagnose:
    """Deck-level diagnostics aggregate per-slide findings and deck-wide issues."""

    def test_should_tag_slide_findings_with_slide_index(self):
        """Per-slide diagnostics carry the originating slide index."""
        # given a slide whose text runs off the canvas
        offending = Canvas(1280, 720).text(
            content="hi", size=64, color="#FFFFFF", position=("200%", "50%")
        )
        deck = Deck(slides=[make_slide("ok"), offending])

        # when
        findings = deck.diagnose()

        # then
        off_canvas = [f for f in findings if f.code == "off-canvas"]
        assert off_canvas
        assert all(f.slide_index == 1 for f in off_canvas)

    def test_should_warn_when_slides_have_mixed_sizes(self):
        """Differing slide dimensions raise a single deck-wide warning."""
        # given
        deck = Deck(slides=[make_slide("a", 1280, 720), make_slide("b", 800, 800)])

        # when
        findings = deck.diagnose()

        # then
        mixed = [f for f in findings if f.code == "mixed-slide-size"]
        assert len(mixed) == 1
        assert mixed[0].slide_index is None
        assert mixed[0].severity == "warning"

    def test_should_not_warn_when_slides_share_a_size(self):
        """Uniformly sized slides produce no mixed-size warning."""
        # given
        deck = Deck(slides=[make_slide("a"), make_slide("b")])

        # when
        findings = deck.diagnose()

        # then
        assert not [f for f in findings if f.code == "mixed-slide-size"]


class TestJsonRoundTrip:
    """Decks serialize to and from JSON via the underlying canvas specs."""

    def test_should_round_trip_through_json(self):
        """from_json(to_json()) reproduces every slide's spec."""
        # given
        deck = Deck(slides=[make_slide("1"), make_slide("2")])

        # when
        restored = Deck.from_json(deck.to_json())

        # then
        assert len(restored) == 2
        assert [c.to_json() for c in restored] == [c.to_json() for c in deck]

    def test_should_round_trip_per_slide_audio_metadata(self):
        """Per-slide narration paths and durations survive Deck JSON serialization"""
        # given
        deck = Deck().slide(make_slide("1"), audio="voice.wav", duration=1.5).slide(make_slide("2"))

        # when
        payload = json.loads(deck.to_json())
        restored = Deck.from_json(deck.to_json())

        # then
        assert payload["kind"] == "deck"
        assert payload["slides"][0]["audio"] == {
            "path": "voice.wav",
            "volume": 1.0,
            "loop": False,
        }
        assert payload["slides"][0]["duration"] == 1.5
        assert json.loads(restored.to_json())["slides"] == payload["slides"]

    def test_should_round_trip_presenter_notes_without_exposing_internal_lists(self):
        """Speaker notes survive JSON while the public notes collection stays isolated."""
        # given: a deck with notes on only its first slide
        deck = Deck().slide(make_slide("1"), notes="Introduce the thesis.").slide(make_slide("2"))

        # when: the deck is restored and the returned notes list is mutated
        restored = Deck.from_json(deck.to_json())
        public_notes = restored.notes
        public_notes[0] = "changed"

        # then: serialized notes survive and callers cannot mutate deck state through the property
        assert restored.notes == ["Introduce the thesis.", None]
        assert json.loads(restored.to_json())["slides"][0]["notes"] == "Introduce the thesis."

    def test_should_reject_non_string_presenter_notes(self):
        """Deck JSON rejects speaker notes that cannot be rendered as text."""
        # given: a deck-shaped JSON document with structured notes
        spec = json.dumps(
            {
                "slides": [
                    {
                        "width": 1280,
                        "height": 720,
                        "layers": [],
                        "notes": ["not", "text"],
                    }
                ]
            }
        )

        # when / then: parsing fails at the public Deck boundary
        with pytest.raises(ValidationError, match="notes"):
            Deck.from_json(spec)

    def test_should_reject_json_without_slides(self):
        """Deck JSON must carry a 'slides' list."""
        # when / then
        with pytest.raises(ValidationError, match="slides"):
            Deck.from_json('{"pages": []}')

    def test_should_preserve_default_size_across_json(self):
        """A restored deck keeps its default size so bare slides still inherit it."""
        # given a sized deck round-tripped through JSON
        deck = Deck(1280, 720, slides=[make_slide("1")])
        restored = Deck.from_json(deck.to_json())

        # when a bare, unsized canvas is added to the restored deck
        restored.slide(Canvas().background(color="#101820"))

        # then it inherits the deck default size instead of raising
        assert (restored[-1].width, restored[-1].height) == (1280, 720)

    def test_should_share_a_deck_level_theme_with_each_slide(self):
        """A top-level theme resolves $theme.* tokens inside every slide."""
        # given a deck spec with a shared theme and a slide that references it
        spec = json.dumps(
            {
                "width": 1280,
                "height": 720,
                "theme": {"brand": "#B8FF00"},
                "slides": [
                    {
                        "width": 1280,
                        "height": 720,
                        "layers": [{"type": "background", "color": "$theme.brand"}],
                    }
                ],
            }
        )

        # when
        deck = Deck.from_json(spec)

        # then the slide's token resolved to the deck-level theme value
        assert cast(BackgroundLayer, deck[0].layers[0]).color == "#B8FF00"
        # and the theme survives another round-trip
        assert json.loads(deck.to_json())["theme"] == {"brand": "#B8FF00"}


P14_NS = "{http://schemas.microsoft.com/office/powerpoint/2010/main}"


class TestDeckTransitions:
    """Deck-level default transitions and per-slide overrides (PPTX only)."""

    @staticmethod
    def _transition_elements(deck: Deck) -> list:
        """The <p:transition> element of each slide, or None when absent."""
        from pptx.oxml.ns import qn

        presentation = Presentation(BytesIO(deck.to_pptx()))
        return [slide._element.find(qn("p:transition")) for slide in presentation.slides]

    @classmethod
    def _transition_effects(cls, deck: Deck) -> list:
        """The transition child tag name on each slide, or None when absent."""
        effects = []
        for element in cls._transition_elements(deck):
            if element is None:
                effects.append(None)
                continue
            children = list(element)
            effects.append(children[0].tag.split("}")[-1] if children else "none")
        return effects

    def test_should_apply_deck_default_transition_to_every_slide(self):
        """A deck default transition is emitted on each slide that lacks its own"""
        # given
        deck = (
            Deck(1280, 720, transition={"effect": "fade"})
            .slide(make_slide("1"))
            .slide(make_slide("2"))
        )

        # when
        effects = self._transition_effects(deck)

        # then
        assert effects == ["fade", "fade"]

    def test_should_let_a_slide_override_win_over_the_deck_default(self):
        """A per-slide transition override wins over the deck default"""
        # given
        deck = Deck(1280, 720, transition={"effect": "fade"})
        deck.slide(make_slide("1"), transition="push")
        deck.slide(make_slide("2"))

        # when
        effects = self._transition_effects(deck)

        # then
        assert effects == ["push", "fade"]

    def test_should_set_transition_inline_via_slide_without_a_default(self):
        """slide(canvas, transition=...) sets the transition for that slide only"""
        # given
        deck = Deck(1280, 720)
        deck.slide(make_slide("1"), transition="wipe")
        deck.slide(make_slide("2"))

        # when
        effects = self._transition_effects(deck)

        # then
        assert effects == ["wipe", None]

    def test_should_emit_duration_speed_and_precise_duration(self):
        """A transition records a speed bucket plus an exact p14 duration"""
        # given
        deck = Deck(1280, 720).slide(
            make_slide("1"), transition={"effect": "fade", "duration": 0.7}
        )

        # when
        element = self._transition_elements(deck)[0]

        # then
        assert element.get("spd") == "med"
        assert element.get(P14_NS + "dur") == "700"

    def test_should_map_direction_to_drawingml_attribute(self):
        """A directional transition writes the matching dir attribute"""
        # given
        from pptx.oxml.ns import qn

        deck = Deck(1280, 720).slide(
            make_slide("1"), transition={"effect": "push", "direction": "left"}
        )

        # when
        push = self._transition_elements(deck)[0].find(qn("p:push"))

        # then
        assert push.get("dir") == "l"

    def test_should_emit_auto_advance_time(self):
        """advance_after writes advTm in milliseconds and keeps advance_on_click"""
        # given
        deck = Deck(1280, 720).slide(
            make_slide("1"),
            transition={"effect": "wipe", "advance_on_click": False, "advance_after": 2.5},
        )

        # when
        element = self._transition_elements(deck)[0]

        # then
        assert element.get("advTm") == "2500"
        assert element.get("advClick") == "0"

    def test_should_accept_a_prebuilt_transition_object(self):
        """A transition effect object can be used as the deck default"""
        # given
        from pptx.oxml.ns import qn

        deck = Deck(1280, 720, transition=Dissolve(duration=1.5)).slide(make_slide("1"))

        # when
        element = self._transition_elements(deck)[0]

        # then
        assert element.find(qn("p:dissolve")) is not None
        assert element.get("spd") == "slow"

    def test_should_emit_effect_specific_transition_options(self):
        """Each transition class emits its own options (spokes, orientation)"""
        # given
        from pptx.oxml.ns import qn

        deck = (
            Deck(1280, 720)
            .slide(make_slide("1"), transition=Wheel(spokes=4))
            .slide(make_slide("2"), transition=Split(orientation="vertical", direction="in"))
        )

        # when
        wheel = self._transition_elements(deck)[0].find(qn("p:wheel"))
        split = self._transition_elements(deck)[1].find(qn("p:split"))

        # then
        assert wheel.get("spokes") == "4"
        assert split.get("orient") == "vert"
        assert split.get("dir") == "in"

    def test_should_reject_non_positive_transition_duration(self):
        """A zero or negative transition duration is a validation error"""
        # given / when / then
        with pytest.raises(ValidationError):
            Fade(duration=0)

    def test_should_raise_a_clean_error_for_an_unknown_transition_effect(self):
        """An unknown effect raises a quickthumb error naming the valid effects"""
        # given / when / then
        with pytest.raises(ValidationError, match="Invalid transition"):
            Deck(1280, 720).slide(make_slide("1"), transition={"effect": "swirl"})

    def test_should_round_trip_effect_specific_transition_fields_through_json(self):
        """Effect-specific transition options survive to_json/from_json"""
        # given
        deck = Deck(1280, 720, transition=Wheel(spokes=4))
        deck.slide(make_slide("1"))
        deck.slide(make_slide("2"), transition=Split(orientation="vertical", direction="in"))

        # when
        restored = Deck.from_json(deck.to_json())

        # then
        default_transition = cast(Wheel, restored.default_transition)
        assert default_transition.spokes == 4
        override = cast(Split, restored._slide_transitions[1])
        assert (override.effect, override.orientation, override.direction) == (
            "split",
            "vertical",
            "in",
        )

    def test_should_round_trip_default_and_per_slide_transitions_through_json(self):
        """Both the deck default and per-slide overrides survive to_json/from_json"""
        # given
        deck = Deck(1280, 720, transition={"effect": "cover", "direction": "up"})
        deck.slide(make_slide("1"))  # uses the default
        deck.slide(make_slide("2"), transition="wipe")  # overrides it

        # when
        restored = Deck.from_json(deck.to_json())

        # then
        assert restored.default_transition is not None
        assert restored.default_transition.effect == "cover"
        assert self._transition_effects(restored) == ["cover", "wipe"]

    def test_should_report_consecutive_identical_transitions(self):
        """Consecutive identical transitions produce an explicit rhythm warning"""
        # given: two adjacent slides with the same transition configuration
        deck = Deck(1280, 720).slide(make_slide("1"), transition=Fade()).slide(
            make_slide("2"), transition=Fade()
        )

        # when
        findings = deck.diagnose()

        # then
        repetition = [finding for finding in findings if finding.code == "transition-repetition"]
        assert len(repetition) == 1
        assert repetition[0].slide_index == 1
