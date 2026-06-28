"""Tests for PPTX export (Canvas.to_pptx and rendering to .pptx files)"""

from io import BytesIO
from pathlib import Path

import pytest
from quickthumb import Animation, Canvas, LinearGradient, TextPart
from quickthumb.errors import RenderingError
from quickthumb.models import Background, Glow, RadialGradient, Shadow, Stroke

pptx = pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from pptx.presentation import Presentation as PresentationDocument  # noqa: E402
from pptx.util import Emu  # noqa: E402

EMU_PER_PX = 9525
FIXTURES_DIR = Path(__file__).parent / "fixtures"
SAMPLE_IMAGE = str(FIXTURES_DIR / "sample_image.jpg")


def timing_of(canvas: Canvas):
    return open_pptx(canvas).slides[0]._element.find(qn("p:timing"))


def open_pptx(canvas: Canvas) -> PresentationDocument:
    return Presentation(BytesIO(canvas.to_pptx()))


def slide_of(canvas: Canvas):
    return open_pptx(canvas).slides[0]


class TestPptxDocument:
    """Test suite for presentation-level output"""

    def test_should_size_slide_to_canvas_pixels(self):
        """The slide adopts the canvas dimensions at 96 dpi"""
        # given
        canvas = Canvas(1280, 720).background(color="#000000")

        # when
        presentation = open_pptx(canvas)

        # then
        assert presentation.slide_width == Emu(1280 * EMU_PER_PX)
        assert presentation.slide_height == Emu(720 * EMU_PER_PX)
        assert len(presentation.slides) == 1

    def test_should_render_pptx_file_from_extension(self, tmp_path):
        """render() writes a PowerPoint file when the output path ends in .pptx"""
        # given
        canvas = Canvas(400, 300).background(color="#FF0000")
        output = tmp_path / "deck.pptx"

        # when
        canvas.render(str(output))

        # then
        assert len(Presentation(str(output)).slides) == 1

    def test_should_raise_rendering_error_for_pptx_canvas_below_slide_minimum(self):
        """PPTX export reports a quickthumb error for canvases below PowerPoint limits"""
        # given
        canvas = Canvas(120, 80)

        # when / then
        with pytest.raises(RenderingError, match="PPTX export supports canvas dimensions"):
            canvas.to_pptx()


class TestPptxBackgroundsAndOutline:
    """Test suite for background and outline layers"""

    def test_should_emit_solid_background_as_full_bleed_rectangle(self):
        """A solid background becomes a full-slide rectangle with the fill color"""
        # given
        canvas = Canvas(400, 300).background(color="#FF8800")

        # when
        slide = slide_of(canvas)

        # then
        shape = slide.shapes[0]
        assert shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        assert shape.width == Emu(400 * EMU_PER_PX)
        assert str(shape.fill.fore_color.rgb) == "FF8800"

    def test_should_emit_linear_gradient_background_as_gradient_fill(self):
        """Linear gradient backgrounds use a native gradient fill"""
        # given
        from pptx.enum.dml import MSO_FILL

        canvas = Canvas(400, 300).background(
            gradient=LinearGradient(angle=90, stops=[("#000000", 0.0), ("#FFFFFF", 1.0)])
        )

        # when
        slide = slide_of(canvas)

        # then
        assert slide.shapes[0].fill.type == MSO_FILL.GRADIENT

    def test_should_emit_outline_as_bordered_rectangle(self):
        """Outline layers become an unfilled rectangle with a matching line"""
        # given
        from pptx.enum.dml import MSO_FILL

        canvas = Canvas(400, 300).outline(width=6, color="#FFFFFF", offset=10)

        # when
        slide = slide_of(canvas)

        # then
        shape = slide.shapes[0]
        assert shape.fill.type == MSO_FILL.BACKGROUND
        assert shape.line.width == Emu(6 * EMU_PER_PX)
        assert str(shape.line.color.rgb) == "FFFFFF"


class TestPptxShapes:
    """Test suite for shape layer mapping"""

    def test_should_map_shape_primitives_to_autoshapes(self):
        """Rectangle, pill, star, and triangle map to the matching autoshapes"""
        # given
        canvas = (
            Canvas(800, 600)
            .shape(
                shape="rectangle",
                position=(10, 10),
                width=100,
                height=60,
                color="#FF0000",
                border_radius=12,
            )
            .shape(shape="pill", position=(150, 10), width=100, height=40, color="#00FF00")
            .shape(
                shape="star",
                position=(300, 10),
                width=80,
                height=80,
                color="#0000FF",
                star_points=5,
            )
            .shape(shape="triangle", position=(420, 10), width=80, height=80, color="#FFFF00")
        )

        # when
        shapes = list(slide_of(canvas).shapes)

        # then
        assert shapes[0].auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
        assert shapes[1].auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
        assert shapes[1].adjustments[0] == pytest.approx(0.5, abs=0.01)
        assert shapes[2].auto_shape_type == MSO_SHAPE.STAR_5_POINT
        assert shapes[3].auto_shape_type == MSO_SHAPE.ISOSCELES_TRIANGLE

    def test_should_emit_polygon_as_freeform(self):
        """Custom polygons become freeform shapes"""
        # given
        canvas = Canvas(400, 300).shape(
            shape="polygon",
            position=(50, 50),
            width=120,
            height=120,
            color="#14B8A6",
            points=[(0.5, 0.0), (1.0, 1.0), (0.0, 1.0)],
        )

        # when
        shape = slide_of(canvas).shapes[0]

        # then
        assert shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
        assert str(shape.fill.fore_color.rgb) == "14B8A6"

    def test_should_apply_rotation_and_stroke_to_shapes(self):
        """Rotation maps to shape.rotation and stroke effects to the shape line"""
        # given
        canvas = Canvas(400, 300).shape(
            shape="rectangle",
            position=(100, 100),
            width=100,
            height=60,
            color="#FFFFFF",
            rotation=15,
            effects=[Stroke(width=3, color="#FF00FF")],
        )

        # when
        shape = slide_of(canvas).shapes[0]

        # then
        assert shape.rotation == pytest.approx(15)
        assert shape.line.width == Emu(3 * EMU_PER_PX)
        assert str(shape.line.color.rgb) == "FF00FF"

    def test_should_apply_shape_opacity_to_stroke_and_effects(self):
        """Shape opacity fades native PPTX stroke, glow, and shadow output"""
        # given
        canvas = Canvas(400, 300).shape(
            shape="rectangle",
            position=(100, 100),
            width=100,
            height=60,
            color="#FFFFFF",
            opacity=0.25,
            effects=[
                Stroke(width=3, color="#FF00FF"),
                Glow(color="#00FFFF", radius=4),
                Shadow(offset_x=4, offset_y=4, color="#000000", blur_radius=2),
            ],
        )

        # when
        shape_xml = slide_of(canvas).shapes[0]._element.xml

        # then
        assert shape_xml.count("<a:alpha") >= 4


class TestPptxText:
    """Test suite for text layer export"""

    def test_should_emit_editable_textbox_with_font_properties(self):
        """Text layers become text boxes with size, weight, color, and font name"""
        # given
        canvas = Canvas(800, 400).text(
            content="Editable headline",
            font="Roboto",
            size=48,
            bold=True,
            color="#FFCC00",
            position=(60, 80),
        )

        # when
        slide = slide_of(canvas)

        # then
        boxes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
        assert len(boxes) == 1
        run = boxes[0].text_frame.paragraphs[0].runs[0]
        assert run.text == "Editable headline"
        assert run.font.size.pt == pytest.approx(48 * 0.75)
        assert run.font.bold is True
        assert str(run.font.color.rgb) == "FFCC00"
        assert run.font.name == "Roboto"

    def test_should_emit_wrapped_text_as_one_paragraph_per_line(self):
        """max_width wrapping becomes explicit paragraphs matching the PNG layout"""
        # given
        canvas = Canvas(400, 300).text(
            content="several words that will certainly wrap around",
            font="Roboto",
            size=40,
            color="#FFFFFF",
            position=(20, 20),
            max_width=300,
        )

        # when
        box = slide_of(canvas).shapes[0]

        # then
        paragraphs = box.text_frame.paragraphs
        assert len(paragraphs) > 1
        rejoined = " ".join(p.runs[0].text for p in paragraphs)
        assert rejoined == "several words that will certainly wrap around"

    def test_should_emit_rich_parts_as_styled_runs(self):
        """Rich content keeps per-part styling within a single paragraph"""
        # given
        canvas = Canvas(600, 200).text(
            content=[
                TextPart(text="Big ", color="#FF0000", size=48, bold=True),
                TextPart(text="small", color="#00FF00", size=24, italic=True),
            ],
            font="Roboto",
            position=(40, 60),
        )

        # when
        box = slide_of(canvas).shapes[0]

        # then
        runs = box.text_frame.paragraphs[0].runs
        assert [run.text for run in runs] == ["Big ", "small"]
        assert runs[0].font.size.pt == pytest.approx(36)
        assert runs[0].font.bold is True
        assert runs[1].font.italic is True
        assert str(runs[1].font.color.rgb) == "00FF00"

    def test_should_emit_valid_drawingml_for_gradient_plus_stroke_run(self):
        """A run with both a gradient fill and a stroke keeps valid rPr child order"""
        # given a gradient-filled headline that also has a stroke effect
        from pptx.oxml.ns import qn

        canvas = Canvas(600, 200).text(
            content="GRAD",
            font="Roboto",
            size=48,
            bold=True,
            position=(40, 60),
            fill=LinearGradient(angle=0, stops=[("#F59E0B", 0.0), ("#EF4444", 1.0)]),
            effects=[Stroke(width=3, color="#7C3AED")],
        )

        # when the exported deck is reopened
        run = slide_of(canvas).shapes[0].text_frame.paragraphs[0].runs[0]
        rpr = run._r.find(qn("a:rPr"))
        tags = [child.tag.split("}")[-1] for child in rpr]

        # then DrawingML requires line, then fill, then latin within rPr
        assert tags.index("ln") < tags.index("gradFill") < tags.index("latin")

    def test_should_rotate_textbox_for_rotated_text(self):
        """Rotation carries over to the text box"""
        # given
        canvas = Canvas(600, 200).text(
            content="tilted",
            font="Roboto",
            size=30,
            color="#FFFFFF",
            position=(300, 100),
            align="center",
            rotation=-20,
        )

        # when
        box = slide_of(canvas).shapes[0]

        # then PowerPoint stores rotation as a positive clockwise angle
        assert box.rotation == pytest.approx(340)

    def test_should_apply_text_opacity_to_background_and_effects(self):
        """Text opacity fades native PPTX backgrounds, strokes, glows, and shadows"""
        # given
        canvas = Canvas(600, 200).text(
            content="badge",
            font="Roboto",
            size=32,
            color="#FFFFFF",
            position=(40, 60),
            opacity=0.25,
            effects=[
                Background(color="#FF0000", padding=8),
                Stroke(width=2, color="#00FF00"),
                Glow(color="#00FFFF", radius=4),
                Shadow(offset_x=2, offset_y=2, color="#000000", blur_radius=2),
            ],
        )

        # when
        slide = slide_of(canvas)
        combined_xml = "\n".join(shape._element.xml for shape in slide.shapes)

        # then
        assert combined_xml.count("<a:alpha") >= 5

    def test_should_rasterize_radial_gradient_text_fill(self):
        """Radial text fills fall back to a picture instead of a flat solid color"""
        # given
        canvas = Canvas(600, 200).text(
            content="RADIAL",
            font="Roboto",
            size=48,
            position=(40, 60),
            fill=RadialGradient(stops=[("#FFFFFF", 0.0), ("#000000", 1.0)]),
        )

        # when
        shape = slide_of(canvas).shapes[0]

        # then
        assert shape.shape_type == MSO_SHAPE_TYPE.PICTURE


class TestPptxEmbeddedLayers:
    """Test suite for layers exported as pictures"""

    def test_should_embed_image_layer_as_picture(self):
        """Image layers become positioned pictures"""
        # given
        canvas = Canvas(400, 300).image(path=SAMPLE_IMAGE, position=(50, 50), width=120)

        # when
        shapes = list(slide_of(canvas).shapes)

        # then
        assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE
        assert shapes[0].left == Emu(50 * EMU_PER_PX)
        assert shapes[0].width == Emu(120 * EMU_PER_PX)

    def test_should_flatten_blend_mode_stack_into_single_picture(self):
        """Blend modes collapse everything beneath them into one picture"""
        # given
        canvas = (
            Canvas(400, 300)
            .background(color="#888888")
            .image(path=SAMPLE_IMAGE, position=(0, 0), width=200, blend_mode="multiply")
            .shape(shape="rectangle", position=(10, 10), width=50, height=50, color="#FF0000")
        )

        # when
        shapes = list(slide_of(canvas).shapes)

        # then
        assert [s.shape_type for s in shapes] == [
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.AUTO_SHAPE,
        ]


class TestPptxElementAnimations:
    """Test suite for per-layer animations emitted as a <p:timing> tree."""

    def test_should_not_emit_timing_without_animations(self):
        """A slide with no animated layers produces no <p:timing> element"""
        # given
        canvas = Canvas(400, 300).background(color="#000000").text(
            "Hi", size=40, position=(10, 10)
        )

        # when
        timing = timing_of(canvas)

        # then
        assert timing is None

    def test_should_emit_entrance_animation_targeting_the_shape(self):
        """An animated shape gets a clickEffect targeting its shape id"""
        # given
        canvas = Canvas(400, 300).shape(
            shape="rectangle",
            position=(10, 10),
            width=80,
            height=40,
            color="#FF0000",
            animation=Animation(effect="fade"),
        )
        document = open_pptx(canvas)
        shape = document.slides[0].shapes[0]

        # when
        timing = document.slides[0]._element.find(qn("p:timing"))
        targets = {t.get("spid") for t in timing.iter(qn("p:spTgt"))}
        node_types = [c.get("nodeType") for c in timing.iter(qn("p:cTn")) if c.get("nodeType")]

        # then
        assert str(shape.shape_id) in targets
        assert "mainSeq" in node_types
        assert "clickEffect" in node_types
        assert [e.get("filter") for e in timing.iter(qn("p:animEffect"))] == ["fade"]

    def test_should_map_wipe_direction_to_animation_filter(self):
        """A directional wipe animation encodes the direction in its filter"""
        # given
        canvas = Canvas(400, 300).text(
            "Slide in",
            size=40,
            position=(10, 10),
            animation=Animation(effect="wipe", direction="right"),
        )

        # when
        timing = timing_of(canvas)

        # then
        assert [e.get("filter") for e in timing.iter(qn("p:animEffect"))] == ["wipe(right)"]

    def test_should_group_with_previous_into_one_click(self):
        """with_previous shares a click group; on_click starts a new one"""
        # given
        canvas = (
            Canvas(400, 300)
            .shape(
                shape="rectangle",
                position=(10, 10),
                width=40,
                height=40,
                color="#FF0000",
                animation=Animation(effect="fade", trigger="on_click"),
            )
            .shape(
                shape="ellipse",
                position=(60, 10),
                width=40,
                height=40,
                color="#00FF00",
                animation=Animation(effect="fade", trigger="with_previous"),
            )
            .shape(
                shape="pill",
                position=(110, 10),
                width=40,
                height=40,
                color="#0000FF",
                animation=Animation(effect="fade", trigger="on_click"),
            )
        )

        # when
        timing = timing_of(canvas)
        main_seq = next(c for c in timing.iter(qn("p:cTn")) if c.get("nodeType") == "mainSeq")
        click_groups = main_seq.find(qn("p:childTnLst")).findall(qn("p:par"))
        node_types = [c.get("nodeType") for c in timing.iter(qn("p:cTn")) if c.get("nodeType")]

        # then
        assert len(click_groups) == 2  # two clicks: [fade+fade], [fade]
        assert node_types.count("clickEffect") == 2
        assert node_types.count("withEffect") == 1

    def test_should_emit_after_previous_as_auto_advancing_group(self):
        """after_previous starts its group automatically (delay 0, afterEffect)"""
        # given
        canvas = (
            Canvas(400, 300)
            .text(
                "First",
                size=40,
                position=(10, 10),
                animation=Animation(effect="fade", trigger="on_click"),
            )
            .text(
                "Second",
                size=40,
                position=(10, 80),
                animation=Animation(effect="fade", trigger="after_previous"),
            )
        )

        # when
        timing = timing_of(canvas)
        node_types = [c.get("nodeType") for c in timing.iter(qn("p:cTn")) if c.get("nodeType")]

        # then
        assert "afterEffect" in node_types

    def test_should_emit_exit_animation_hiding_the_shape(self):
        """An exit animation plays out and pins visibility to hidden"""
        # given
        canvas = Canvas(400, 300).shape(
            shape="rectangle",
            position=(10, 10),
            width=80,
            height=40,
            color="#FF0000",
            animation=Animation(effect="fade", animate="exit"),
        )

        # when
        timing = timing_of(canvas)
        effect = next(c for c in timing.iter(qn("p:cTn")) if c.get("presetClass"))
        visibility = [v.get("val") for v in timing.iter(qn("p:strVal"))]

        # then
        assert effect.get("presetClass") == "exit"
        assert "hidden" in visibility

    def test_should_animate_all_shapes_a_text_layer_produces(self):
        """A text layer with a background fill animates the box and its background together"""
        # given
        canvas = Canvas(400, 300).text(
            "Boxed",
            size=40,
            position=(10, 10),
            effects=[Background(color="#222222", padding=8)],
            animation=Animation(effect="fade"),
        )
        document = open_pptx(canvas)

        # when
        timing = document.slides[0]._element.find(qn("p:timing"))
        animated = {t.get("spid") for t in timing.iter(qn("p:spTgt"))}
        all_ids = {str(s.shape_id) for s in document.slides[0].shapes}

        # then: every shape the layer produced is animated (background + text box)
        assert len(animated) >= 2
        assert animated <= all_ids

    def test_should_animate_a_group_as_a_single_effect(self):
        """A group animation plays across all its children on one click"""
        # given
        canvas = Canvas(1280, 720).group(
            children=[
                {"type": "text", "content": "Agenda", "size": 90, "color": "#FFFFFF"},
                {"type": "shape", "shape": "pill", "width": 200, "height": 40, "color": "#B8FF00"},
            ],
            position=("50%", "50%"),
            align="center",
            animation=Animation(effect="fade", trigger="on_click"),
        )

        # when
        timing = timing_of(canvas)
        main_seq = next(c for c in timing.iter(qn("p:cTn")) if c.get("nodeType") == "mainSeq")
        click_groups = main_seq.find(qn("p:childTnLst")).findall(qn("p:par"))
        node_types = [c.get("nodeType") for c in timing.iter(qn("p:cTn")) if c.get("nodeType")]
        animated = {t.get("spid") for t in timing.iter(qn("p:spTgt"))}

        # then: both children animate, together, under a single click
        assert len(click_groups) == 1
        assert node_types.count("clickEffect") == 1
        assert node_types.count("withEffect") == 1
        assert len(animated) == 2

    def test_group_animation_should_take_precedence_over_child_animations(self):
        """A group animation drives all children as one effect, ignoring child ones"""
        # given a group animation plus a conflicting child animation
        canvas = Canvas(1280, 720).group(
            children=[
                {"type": "text", "content": "A", "size": 60, "color": "#FFFFFF"},
                {
                    "type": "shape",
                    "shape": "pill",
                    "width": 120,
                    "height": 40,
                    "color": "#B8FF00",
                    "animation": {"effect": "wipe", "trigger": "on_click"},
                },
            ],
            position=("50%", "50%"),
            align="center",
            animation=Animation(effect="fade", trigger="on_click"),
        )

        # when
        timing = timing_of(canvas)
        main_seq = next(c for c in timing.iter(qn("p:cTn")) if c.get("nodeType") == "mainSeq")
        click_groups = main_seq.find(qn("p:childTnLst")).findall(qn("p:par"))
        filters = [e.get("filter") for e in timing.iter(qn("p:animEffect"))]

        # then: one click, both children fade (the child's own wipe is overridden)
        assert len(click_groups) == 1
        assert filters == ["fade", "fade"]

    def test_should_reproduce_animation_through_json_round_trip(self):
        """Animations survive to_json/from_json and still emit timing"""
        # given
        canvas = Canvas(400, 300).text(
            "Hi", size=40, position=(10, 10), animation=Animation(effect="wipe", direction="up")
        )

        # when
        restored = Canvas.from_json(canvas.to_json())
        timing = timing_of(restored)

        # then
        assert timing is not None
        assert [e.get("filter") for e in timing.iter(qn("p:animEffect"))] == ["wipe(up)"]
